#!/usr/bin/env python3
"""
Скрипт для замены текста и изображений в презентации PowerPoint.
Заменяет все вхождения {дата+3} на 21.09.25 и {image0} на изображение
"""

import traceback
from pptx import Presentation
from pptx.util import Cm
import os
from datetime import datetime, timedelta
import pytz


def replace_text_and_images_in_presentation(input_file, output_file, text_replacements, image_replacements):
    """
    Заменяет текст и изображения в презентации PowerPoint
    
    Args:
        input_file (str): Путь к исходному файлу презентации
        output_file (str): Путь к выходному файлу
        text_replacements (dict): Словарь замен текста {старый_текст: новый_текст}
        image_replacements (dict): Словарь замен изображений {текст_заглушка: (путь_к_изображению, ширина_см, высота_см)}
    """
    try:
        # Загружаем презентацию
        prs = Presentation(input_file)
        
        # Счетчики замен
        text_replacements_count = 0
        image_replacements_count = 0
        
        # Проходим по всем слайдам
        for slide_num, slide in enumerate(prs.slides, 1):
            print(f"Обрабатываем слайд {slide_num}...")
            
            # Проходим по всем фигурам на слайде
            for shape in slide.shapes:
                # Обработка текстовых замен
                if shape.has_text_frame:
                    # Проходим по всем параграфам в текстовом фрейме
                    for paragraph in shape.text_frame.paragraphs:
                        # Проходим по всем текстовым блокам (runs) в параграфе
                        for run in paragraph.runs:
                            # Заменяем текст
                            for old_text, new_text in text_replacements.items():
                                if old_text in run.text:
                                    print(f"  Заменяем текст: {old_text} -> {new_text}")
                                    run.text = run.text.replace(old_text, new_text)
                                    text_replacements_count += 1
                                    print(f"  Заменен текст в фигуре: {shape.name}")
                
                # Обработка замены изображений
                for placeholder_text, image_data in image_replacements.items():
                    if shape.has_text_frame and placeholder_text in shape.text:
                        # Если значение пустое - просто удаляем текст
                        if not image_data or image_data == '':
                            shape.text = shape.text.replace(placeholder_text, "")
                            print(f"  Удален текст-заглушка: {placeholder_text}")
                            continue
                        
                        # Распаковываем данные изображения
                        image_path, width_cm, height_cm = image_data
                        
                        # Проверяем существование файла изображения
                        if not os.path.exists(image_path):
                            print(f"  Предупреждение: Файл изображения {image_path} не найден!")
                            continue
                        
                        # Удаляем текстовую заглушку
                        shape.text = shape.text.replace(placeholder_text, "")
                        
                        # Добавляем изображение
                        try:
                            # Получаем позицию и размеры фигуры
                            left = shape.left
                            top = shape.top
                            
                            # Добавляем изображение с указанными размерами
                            pic = slide.shapes.add_picture(
                                image_path, 
                                left, 
                                top, 
                                width=Cm(width_cm), 
                                height=Cm(height_cm)
                            )
                            
                            image_replacements_count += 1
                            print(f"  Добавлено изображение в фигуру: {shape.name}")
                            
                        except Exception as e:
                            print(f"  Ошибка при добавлении изображения: {traceback.format_exc()}")
        
        # Сохраняем обновленную презентацию
        prs.save(output_file)
        
        print(f"\nОбработка завершена!")
        print(f"Замен текста: {text_replacements_count}")
        print(f"Замен изображений: {image_replacements_count}")
        print(f"Результат сохранен в: {output_file}")
        
        return True
        
    except FileNotFoundError:
        print(f"Ошибка: Файл {input_file} не найден!")
        return False
    except Exception as e:
        print(f"Ошибка при обработке презентации: {traceback.format_exc()}")
        return False


def create_presentation(frakcia,ypakovka,dostavka,opportunity,productName,images,productPrice,obem_po_porametram):
    """Основная функция"""
    import os
    # Пути к файлам
    input_file = "КП - итоговый.pptx"
    output_file = "кп 3 вариант (1)_обновленный.pptx"
    
    tz=pytz.timezone('Europe/Moscow')
    date = datetime.now(tz) + timedelta(days=4)
    # Настройки замен текста
    # Иногда productName приходит с переносами строк и лишними пробелами, убираем их
    # Исправляем: убираем только переносы строк, но сохраняем пробелы между словами
    productName = ' '.join(productName.replace('\n', ' ').replace('\r', ' ').split())
    text_replacements = {
        "{дата+3}": date.strftime("%d.%m.%y"),
        '{Название товара}': productName,
        "{fract}": frakcia,
        "{Объем и вес}": obem_po_porametram,
        "{Цена}": str(float(productPrice))+ " р.",
        "{Стоимость}": str(float(opportunity))+ " р.",
        "{Доставка}": "с учетом доставки" if dostavka else "без учета доставки",
        "{Упаковка}": ypakovka,
       
    }
    
    # Настройки замен изображений
    image_replacements = {
        "{image0}": (images['default'], 3, 3) if images['default'] else '',
        "{image1}": (images['dry'], 3, 3) if images['dry'] else '',
        "{image2}": (images['wet'], 3, 3) if images['wet'] else '',
        "{image3}": (images['lit'], 3, 3) if images['lit'] else ''
    }
    
    print("=" * 60)
    print("ЗАМЕНА ТЕКСТА И ИЗОБРАЖЕНИЙ В ПРЕЗЕНТАЦИИ")
    print("=" * 60)
    print(f"Исходный файл: {input_file}")
    print(f"Выходной файл: {output_file}")
    print("\nЗамены текста:")
    for old, new in text_replacements.items():
        print(f"  '{old}' -> '{new}'")
    print("\nЗамены изображений:")
    for placeholder, image_data in image_replacements.items():
        if not image_data or image_data == '':
            print(f"  '{placeholder}' -> удалить")
        else:
            path, w, h = image_data
            print(f"  '{placeholder}' -> '{path}' ({w}см x {h}см)")
    print("=" * 60)
    
    # Проверяем существование исходного файла
    if not os.path.exists(input_file):
        print(f"Ошибка: Файл {input_file} не найден в текущей директории!")
        return
    
    # Проверяем существование файлов изображений
    for placeholder, image_data in image_replacements.items():
        # Пропускаем пустые значения
        if not image_data or image_data == '':
            continue
        image_path, _, _ = image_data
        if not os.path.exists(image_path):
            print(f"Ошибка: Файл изображения {image_path} не найден!")
            return
    
    # Выполняем замену
    success = replace_text_and_images_in_presentation(input_file, output_file, text_replacements, image_replacements)
    
    if success:
        print("\n✅ Операция выполнена успешно!")
    else:
        print("\n❌ Операция завершилась с ошибкой!")

    

    input_pptx = "кп 3 вариант (1)_обновленный.pptx"
    output_pdf = f"КП {date.strftime('%d.%m.%y')}_{productName}.pdf"

    # Проверяем существование входного файла
    if not os.path.exists(input_pptx):
        print(f"Ошибка: Файл {input_pptx} не найден!")
        return None

    # Используем LibreOffice headless режим вместо unoconv
    try:
        import subprocess
        # Команда для конвертации через LibreOffice
        command = [
            "libreoffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", os.path.dirname(os.path.abspath(output_pdf)) or ".",
            input_pptx
        ]
        
        result = subprocess.run(command, capture_output=True, text=True, timeout=60)
        
        if result.returncode == 0:
            # LibreOffice создает файл с тем же именем, но с расширением .pdf
            base_name = os.path.splitext(os.path.basename(input_pptx))[0]
            generated_pdf = f"{base_name}.pdf"
            
            # Переименовываем в нужное имя
            if os.path.exists(generated_pdf):
                if generated_pdf != output_pdf:
                    os.rename(generated_pdf, output_pdf)
                print(f"PDF успешно создан: {output_pdf}")
                return output_pdf
            else:
                print(f"Ошибка: PDF файл не был создан")
                return None
        else:
            print(f"Ошибка конвертации: {result.stderr}")
            return None
            
    except subprocess.TimeoutExpired:
        print("Ошибка: Превышено время ожидания конвертации")
        return None
    except FileNotFoundError:
        print("Ошибка: LibreOffice не найден. Установите LibreOffice или используйте альтернативный метод")
        return None
    except Exception as e:
        print(f"Ошибка при конвертации: {e}")
        return None

if __name__ == "__main__":
    create_presentation()
    